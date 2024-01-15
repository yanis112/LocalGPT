import os
import chromadb
from flask import Flask, render_template, request
from tkinter import Tk, filedialog
import fitz  # PyMuPDF
from pdf2image import convert_from_path
from gtts import gTTS
from speach_to_text_whisper import split_audio_and_transcribe
import pytesseract
import email
import imaplib
import numpy as np
import re
from Voice_Assistant_V2 import gpt_answer, record_audio, text_to_speech
from pptx import Presentation
from PIL import Image
from HuggingChat import huggingface_chatbot
from CalendarGestion import get_calendar
from InternetGPT import *
import html
import hashlib



app = Flask(__name__)
UPLOAD_FOLDER = ""** PUT YOUR FOLDER PATH HERE **""


app.config['JSON_AS_ASCII'] = False
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def generer_clef(texte):
    # Utiliser hashlib pour calculer une empreinte unique
    h = hashlib.sha256(texte.encode())
    # Convertir l'empreinte en une représentation hexadécimale
    cle = h.hexdigest()
    return cle

class DirectoryToChroma:
    def __init__(self, collection_name, folder_path, chunk_size):
        self.chroma_client = chromadb.PersistentClient(path="c:/Users/Yanis/Documents/chroma/")
        try:
            self.collection = self.chroma_client.get_collection(name=collection_name)
        except:
            self.collection = self.chroma_client.create_collection(name=collection_name)

        self.folder_path = folder_path
        self.chunk_size = chunk_size

    def _segment_text(self, text):
        return [text[i:i + self.chunk_size] for i in range(0, len(text), self.chunk_size)]
    
    def extract_text_from_pdf(self, pdf_path):
        pdf_text = ""
        try:
            doc = fitz.open(pdf_path)
            num_pages = doc.page_count
            for page_num in range(num_pages):
                print("Starting Ocerising page: "+str(page_num))
                page = doc.load_page(page_num)
                pdf_text += page.get_text()
                
                # Extract images and perform OCR
                images = page.get_images(full=True)
                for img_index, img in enumerate(images):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_data = base_image["image"]
                    image = fitz.Pixmap(doc, xref)
                    
                    if image.n >= 4:  # RGBA or CMYK
                        image = fitz.Pixmap(fitz.csRGB, image)
                    
                    img_path = f"temp_image_{page_num}_{img_index}.png"
                    image.save(img_path, output="png")
                    
                    ocr_text = self.perform_ocr(img_path)
                    pdf_text += ocr_text
                    
                    #image.close()
                    os.remove(img_path)
                
            doc.close()
        except Exception as e:
            print("Error extracting text from PDF:", e)
        return pdf_text

    def perform_ocr(self, image_path):
        try:
            ocr_text = pytesseract.image_to_string(image_path)
            return ocr_text
        except Exception as e:
            print("Error performing OCR:", e)
            return ""

    def process_directory(self, folder_path=None):
        if folder_path is None:
            folder_path = self.folder_path

        for item in os.listdir(folder_path):
            item_path = os.path.join(folder_path, item)
            print(item_path)
            if os.path.isfile(item_path):
                if item.lower().endswith(('.pdf', '.txt', '.py','.pptx')):
                    try:
                        if item.lower().endswith('.pdf'):
                            print("extraction of pdf starting")
                            pdf_text = self.extract_text_from_pdf(item_path)
                            chunks = self._segment_text(pdf_text)
                        elif item.lower().endswith('.pptx'):
                            print("PPTX file DETECTED")
                            presentation = Presentation(item_path)
                            presentation_text = ""

                            for slide in presentation.slides:
                                slide_text = ""
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        slide_text += shape.text + "\n"
                                    elif shape.shape_type == 3:  # Shape type 3 corresponds à une image
                                        # Convertir l'image en PIL Image
                                        img = shape.image
                                        image_bytes = img.blob
                                        with open("temp_image.png", "wb") as img_file:
                                            img_file.write(image_bytes)
                                        image = Image.open("temp_image.png")

                                        # Utiliser pytesseract pour extraire le texte de l'image
                                        print("Ocerisation starting")
                                        image_text = pytesseract.image_to_string(image, lang="eng")

                                        # Ajouter le texte de l'image au texte de la diapositive
                                        slide_text += image_text + "\n"

                                presentation_text += slide_text
                                print(presentation_text)

                            chunks = self._segment_text(presentation_text)
                        else:
                            with open(item_path, 'r', encoding='iso-8859-1') as file:
                                content = file.read()
                            chunks = self._segment_text(content)
                            
                        for i, chunk in enumerate(chunks):
                            identifier = f"{item}_chunk{i}"
                            self.collection.add(
                                documents=[chunk],
                                metadatas=[{"source": item}],
                                ids=[identifier]
                            )
                        print("File " + str(item) + " processed!")
                    except Exception as e:
                        print("Error processing", item, ":", e)

            elif os.path.isdir(item_path):
                print("Subfolder detected:", item)
                self.process_directory(item_path)  # Appel récursif pour traiter les sous-dossiers

        print("Processing and adding chunks to Chroma collection complete.")

    def query(self, query_text, n_results):
        #chroma_client = chromadb.PersistentClient(path="c:/Users/Yanis/Documents/chroma/")
        #collection = chroma_client.get_collection(name=collection_name)
        results = self.collection.query(
            query_texts=[query_text],
            n_results=n_results
        )
        #print(results)
        res=str(results['documents'][0])
        #res=results['documents'][0]
        prompt=str(query_text)+", Répond en Français et en moins de 100 mots et bases toi sur le contexte suivant entre étoiles: "+ "** "+ str(res)+ " ** " +". Pour répondre, si tu ne trouve pas l'information demandée dans le contexte, répond: **je ne sais pas**."
        print(prompt)
        answer=huggingface_chatbot(prompt)
        return answer #results['documents'][0],

def index_emails(collection, folder_path):
    email_contents_array = get_emails()
    
    for i, email_content in enumerate(email_contents_array):
        print(email_content)
        print("email "+str(i)+"indexed !")
        cleaned_content = clean_text(email_content)
        identifier = "email_"+str(generer_clef(cleaned_content))              #_{cleaned_content[0:20]}"+str(i)
        try:
            collection.delete(ids=[identifier])
            print('Already exesting mail deleted ')
        except:
            print("This identifier is not in the database")
        collection.add(
            documents=[cleaned_content],
            metadatas=[{"source": "email"}],
            ids=[identifier]
        )
        print(f"Email {i} indexed!")
    
    #calendar_synchro
    calendar=get_calendar()
    id1="calendar1"
    try:
        collection.delete(ids=[id1])
    except:
        print("calendar doesnt exept")
    collection.add(
            documents=[calendar],
            metadatas=[{"source": "email"}],
            ids=[id1]
        )
    print("Calendar Added !")

def get_text_after_newline(text):
    index = text.find("\n")  # Trouver l'index du premier caractère "\n"
    if index != -1:  # Si le caractère "\n" est trouvé
        return text[index+1:]  # Renvoyer la partie du texte après le "\n"
    else:
        return ""  # Si le caractère "\n" n'est pas trouvé, renvoyer une chaîne vide


def get_emails():
    imap_server = 'YOUR IMAP ADRESS'
    email_address = 'YOUR EMAIL ADRESS'
    password='YOUR PASSWORD'

    # Connexion IMAP
    imap = imaplib.IMAP4_SSL(imap_server)
    imap.login(email_address, password)
    imap.select("Inbox")

    # Recherche des 100 derniers e-mails
    _, msgnums = imap.search(None, "ALL")
    latest_msgnums = msgnums[0].split()[-100:]  # Sélectionne les 100 derniers e-mails

    # Liste pour stocker les contenus des e-mails
    email_contents = []
    email_senders = []

    # Récupération du contenu de chaque e-mail
    for msgnum in latest_msgnums:
        _, data = imap.fetch(msgnum, "(RFC822)")
        message = email.message_from_bytes(data[0][1])

        # Récupérer l'expéditeur de l'e-mail à partir de l'en-tête "From"
        sender = message["From"]
        email_senders.append(sender)

        # Ajouter le contenu de l'e-mail à la liste
        for part in message.walk():
            if part.get_content_type() == "text/plain":
                decoded_content = part.get_payload(decode=True).decode('utf-8', errors='ignore')
                email_contents.append(decoded_content)

    imap.close()

    # Convertir la liste en un tableau numpy
    email_contents_array = np.array(email_contents)
    return(email_contents_array)

def special_caracters(texte):
    e=texte.replace("&#39;"," ")
    e=texte.replace("&#34;"," ")
    e=texte.replace("'"," ")
    return(e)


def clean_text(text):
    # Supprimer les caractères spéciaux, sauf les accents
    text = re.sub(r"[^\w\sàâäéèêëîïôöùûüç:/]", "", text)

    # Supprimer les envoyeurs de mails et destinataires
    text = re.sub(r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b", "", text)

    # Supprimer les chaînes de caractères précédées d'un anti slash,
    # sauf celles associées aux accents
    text = re.sub(r"(?<!\\)(\\(?!['`^~]))", " ", text)
    text = re.sub(r"[\r\n\t]", " ", text)

    # Supprimer les espaces inutiles
    text = re.sub(r"\s+", " ", text).strip()
    text = text.replace('\xa0', ' ')

    # Convertir le texte en minuscules
    text = text.lower()
    return text


# Instance de la classe
collection_name = "document_collection"
default_folder_path = "C:/Users/Yanis/Documents/KSI/transcript/"
chunk_size = 1000
directory_handler = None  # L'instance sera créée plus tard

# ... Autres importations ...

@app.route("/", methods=["GET", "POST"])
def index():
   
    global directory_handler
    folder_path = "C:/Users/Yanis/Documents/KSI/empty/"
    directory_handler = DirectoryToChroma(collection_name, folder_path, chunk_size)
    # Initialisez results avec une chaîne vide


    if request.method == "POST":
        folder_path = request.form.get("folder_path")
        query_text = request.form.get("query_text")
        index_emails_flag = request.form.get("index_emails")
        record_audio_flag = request.form.get("record_audio")
        web_query = request.form.get("web_query")

        if folder_path:
            directory_handler = DirectoryToChroma(collection_name, folder_path, chunk_size)
            directory_handler.process_directory()

        if web_query:
            answer=web_GPT(str(web_query))
            final = html.unescape(special_caracters(answer))
            final=re.sub(r"\s+", " ", final).strip()
            return render_template("index.html",results=final)


        if index_emails_flag:
            index_emails(directory_handler.collection, folder_path)  # Indexe les e-mails
    

        if record_audio_flag:
            # Enregistre l'audio
            WAVE_OUTPUT_FILENAME = "output.wav"
            record_audio(WAVE_OUTPUT_FILENAME, 5)

            # Transcription audio en texte
            print("*Transcription starting*")
            texte = split_audio_and_transcribe(WAVE_OUTPUT_FILENAME,app.config['UPLOAD_FOLDER'], 5000)
            print("*Transcription finished*")

            # Query utilisant la transcription comme contexte
            if directory_handler:
                num_results = 3
                query_results = directory_handler.query(texte, num_results)
                context = " ".join([result[0] for result in query_results])

                # Appel à la fonction d'obtention d'une réponse via GPT-3
                print("*Prompt answer loading*")
                few_shot_prompt = "Answer in french: " + context
                answer = gpt_answer(few_shot_prompt + texte)
                print(answer)

                # Conversion du texte en discours audio
                text_to_speech(answer)
                print("*Audio response generated*")

        if query_text and directory_handler:
            num_results = 3
            query_results = directory_handler.query(query_text, num_results)
            final = html.unescape(special_caracters(query_results))
            final=re.sub(r"\s+", " ", final).strip()
            print("QUERY RESULTS:",query_results)
            return render_template("index.html",results= final)

    return render_template("index.html", results=None)  #results=None)

if __name__ == "__main__":
    app.run(debug=True)

